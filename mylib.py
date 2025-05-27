import functools
import hashlib
import json
import os
import re
import ssl
import socket
import subprocess
import time
import urllib
from datetime import datetime
from pathlib import Path

import pandas as pd
import requests
from requests.exceptions import Timeout
from akamai.edgegrid import EdgeGridAuth, EdgeRc
from cryptography import x509

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

pd.options.display.max_rows = None
pd.options.display.max_columns = None

# Prepare API credential
file = os.path.join(Path.home(), ".edgerc")
try:
    EDGE_RC = EdgeRc(file)
    SECTION = "default"
    BASE_URL = f"https://{EDGE_RC.get(SECTION, 'host')}"
    SESSION = requests.Session()
    SESSION.auth = EdgeGridAuth.from_edgerc(EDGE_RC, SECTION)
except:
    raise ValueError(
        f"Create '{file}' by referring https://collaborate.akamai.com/confluence/display/~rosawa/Get+Started+with+Akamai+API"
    )


def try_except(func):
    @functools.wraps(func)
    def wrapper(*args, **kwargs):

        try:
            return func(*args, **kwargs)

        except Exception as e:
            print(f"{e} ({func.__name__})")

    return wrapper


def extract_row_of_dataframe(df, display_table=False):

    df = df.reset_index(drop=True)

    if not len(df):
        raise ValueError("No available data")
    elif (len(df) > 1) or display_table:
        try:
            display(df)
        except:
            print(f"\n{df}")

    if len(df) == 1:
        i = 0
    else:
        i = input("Select index:")
        i = int(i) if i.isdigit() else None

    if type(i) is int:
        return df.loc[i].to_dict()
    else:
        raise ValueError("Invalid index")


def parse_certificate(cert):

    cert = x509.load_pem_x509_certificate(cert.encode())

    result = {}

    result["CN"] = cert.subject.rfc4514_string()
    result["Issuer"] = cert.issuer.rfc4514_string()
    result["Not Before"] = cert.not_valid_before_utc
    result["Not After"] = cert.not_valid_after_utc
    result["Serial"] = format(cert.serial_number, "X")
    result["SAN"] = []

    for ext in cert.extensions:
        if isinstance(ext.value, x509.SubjectAlternativeName):
            for v in ext.value:
                result["SAN"].append(v.value)

    return result


def extract_behaviors(rules, target_behaviors=None):

    # Check criteria
    criteria = rules.get("criteria")

    if criteria:
        criteria_texts = []

        for criterion in criteria:
            # key
            if criterion["options"].get("matchOperator"):
                if criterion["name"] == "matchVariable":
                    text = f"{criterion['options']['variableName']} {criterion['options']['matchOperator'].lower().replace('_',' ')}"
                elif criterion["name"] in ["requestHeader", "responseHeader"]:
                    text = f"{criterion['name']}({criterion['options']['headerName']}) {criterion['options']['matchOperator'].lower().replace('_',' ')}"
                else:
                    text = f"{criterion['name']} {criterion['options']['matchOperator'].lower().replace('_',' ')}"
            else:
                text = criterion["name"]

            # value
            for option in [
                "originType",
                "network",
                "countries",
                "countryValues",
                "variableValues",
                "matchOn",
                "matchString",
                "value",
                "values",
            ]:
                if criterion["options"].get(option):
                    text += f" {criterion['options'][option]}"
                    break

            criteria_texts.append(text)

        if rules["criteriaMustSatisfy"] == "all":
            rules["criteria_text"] = " AND ".join(criteria_texts)
        elif rules["criteriaMustSatisfy"] == "any":
            rules["criteria_text"] = " OR ".join(criteria_texts)
        else:
            raise ValueError(
                f"Unexpected match operator: {rules['criteriaMustSatisfy']}"
            )

        if rules.get("criteria_chain"):
            rules["criteria_chain"].append(rules["criteria_text"])
        else:
            rules["criteria_chain"] = [rules["criteria_text"]]

    else:
        rules["criteria_text"] = None
        rules["criteria_chain"] = (
            rules["criteria_chain"] if rules.get("criteria_chain") else []
        )

    # Check behaviors
    result = []
    rule_name = re.sub(r"^default > ", "", rules["name"])
    criteria = " AND ".join(
        [
            f"({criteria})" if len(rules["criteria_chain"]) > 2 else criteria
            for criteria in rules["criteria_chain"]
        ]
    )

    for behavior in rules["behaviors"]:
        if target_behaviors and (behavior["name"] not in target_behaviors):
            continue
        behavior["rule_name"] = rule_name
        behavior["criteria"] = criteria
        result.append(behavior)

    # Check Child rules
    for child_rule in rules["children"]:
        if rules["name"] not in child_rule["name"]:
            child_rule["name"] = f"{rules['name']} > {child_rule['name']}"
        child_rule["criteria_chain"] = rules["criteria_chain"].copy()
        result.extend(extract_behaviors(child_rule, target_behaviors=target_behaviors))

    return result


def extract_cpcodes(rules):

    behaviors = extract_behaviors(rules, target_behaviors=["cpCode"])
    result = []

    for behavior in behaviors:
        result.append(behavior["options"]["value"].get("id"))

    return result


def extract_im_cpcodes(rules):

    behaviors = extract_behaviors(
        rules, target_behaviors=["imageManager", "imageManagerVideo"]
    )
    result = {"pristine": [], "derivative": []}

    for behavior in behaviors:
        if behavior["options"].get("cpCodeOriginal"):
            result["pristine"].append(behavior["options"]["cpCodeOriginal"].get("id"))
        if behavior["options"].get("cpCodeTransformed"):
            result["derivative"].append(
                behavior["options"]["cpCodeTransformed"].get("id")
            )

    return result


def extract_im_policies(rules):

    behaviors = extract_behaviors(
        rules, target_behaviors=["imageManager", "imageManagerVideo"]
    )
    result = []

    for behavior in behaviors:
        result.append(behavior["options"].get("policyTokenDefault"))

    return result


def extract_ssmaps(rules):

    behaviors = extract_behaviors(rules, target_behaviors=["siteShield"])
    result = []

    for behavior in behaviors:
        result.append(behavior["options"]["ssmap"].get("value"))

    return result


def extract_sr_ssmaps(rules):

    behaviors = extract_behaviors(rules, target_behaviors=["siteShield"])
    result = []

    for behavior in behaviors:
        result.append(behavior["options"]["ssmap"].get("srmap"))

    return result


def extract_srto_url(rules, hostname=None):

    behaviors = extract_behaviors(rules, target_behaviors=["sureRoute"])
    result = []

    for behavior in behaviors:
        if behavior["options"].get("testObjectUrl"):
            result.append(behavior["options"]["testObjectUrl"])

    return result


def extract_cache_key_hostnames(rules, request_host_header):

    behaviors = extract_behaviors(rules, target_behaviors=["origin"])
    result = []

    for behavior in behaviors:
        if behavior["options"].get("cacheKeyHostname") == "REQUEST_HOST_HEADER":
            result.append(request_host_header)
        elif behavior["options"].get("cacheKeyHostname") == "ORIGIN_HOSTNAME":
            result.append(behavior["options"].get("hostname"))
        elif behavior["options"].get("netStorage"):
            result.append(behavior["options"]["netStorage"].get("downloadDomainName"))

    return result


def extract_origin_hostnames(rules):

    behaviors = extract_behaviors(rules, target_behaviors=["origin"])
    result = []

    for behavior in behaviors:
        if behavior["options"].get("hostname"):
            result.append(behavior["options"]["hostname"])
        elif behavior["options"].get("netStorage"):
            result.append(behavior["options"]["netStorage"].get("downloadDomainName"))

    return result


def extract_caching_options(rules):

    behaviors = extract_behaviors(rules, target_behaviors=["caching"])
    result = []

    for behavior in behaviors:
        d = dict()
        d["option"] = [behavior["options"]["behavior"]]

        for key in ["ttl", "defaultTtl"]:
            ttl = behavior["options"].get(key)

            if ttl is not None:
                num = int(re.search(r"[0-9]+", ttl).group())
                unit = {
                    "d": "days",
                    "h": "hours",
                    "s": "seconds",
                    "m": "minutes",
                }.get(re.search(r"[a-z]+", ttl).group())
                d["ttl"] = [num, unit]
                break

        d["criteria"] = behavior["criteria"]

        if d not in result:
            result.append(d)

    return result


def extract_lma_options(rules):

    behaviors = extract_behaviors(rules, target_behaviors=["gzipResponse"])
    result = []

    for behavior in behaviors:
        d = dict()
        d["options"] = behavior["options"]["behavior"]
        d["criteria"] = behavior["criteria"]

        if d not in result:
            result.append(d)

    return result


def extract_enhanced_debug_keys(rules):

    behaviors = extract_behaviors(rules, target_behaviors=["enhancedDebug"])
    result = []

    for behavior in behaviors:
        if behavior["options"].get("debugKey"):
            result.append(behavior["options"]["debugKey"])

    return result


def run_command(cmd, timeout=10):

    result = subprocess.run(
        cmd.split(" "),
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=timeout,
    )

    print(" ".join(result.args))
    print(result.stdout)
    if result.stderr:
        print(result.stderr)

    return result.stdout


def get_cert(hostname, sni=None, port=443):

    context = ssl.create_default_context()
    with socket.create_connection((hostname, port)) as sock:
        with context.wrap_socket(
            sock, server_hostname=sni if sni else hostname
        ) as ssock:
            cert = ssock.getpeercert(binary_form=True)

    cert = ssl.DER_cert_to_PEM_cert(cert)
    cert = parse_certificate(cert)

    return cert


def generate_etag(json_obj):

    json_str = json.dumps(json_obj)
    hash_str = hashlib.md5(json_str.encode()).hexdigest()
    etag = f'"{hash_str}"'

    return etag


def convert_bytes(value):

    value = int(value)
    units = [
        "B",
        "KB",
        "MB",
        "GB",
        "TB",
        "PB",
    ]
    unit_index = 0

    while value >= 1000 and unit_index < len(units) - 1:
        value /= 1000
        unit_index += 1

    converted_value = f"{value:.2f} {units[unit_index]}"

    return converted_value


def get_enhanced_debug_token(key):

    url = f"https://debug.akamai.com/generate/token?key={key}&duration=1d&token_type=auth_token&format=api"
    response = requests.get(url)
    result = response.text
    return result


def access_api(
    method,
    path,
    query,
    header=None,
    payload=None,
    error_threshold=400,
    max_retry=2,
    retry_interval=60,
    timeout=120,
    session=SESSION,
    base_url=BASE_URL,
):

    url = urllib.parse.urljoin(base_url, path)
    headers = {"accept": "application/json"}

    if header:
        headers.update(header)

    retry_counter = 0

    while retry_counter <= max_retry:
        try:
            if method == "GET":
                response = session.get(
                    url, headers=headers, params=query, timeout=timeout
                )

            elif method == "POST":
                response = session.post(
                    url, headers=headers, json=payload, params=query, timeout=timeout
                )

            elif method == "PUT":
                response = session.put(
                    url, headers=headers, json=payload, params=query, timeout=timeout
                )

            elif method == "PATCH":
                response = session.patch(
                    url, headers=headers, json=payload, params=query, timeout=timeout
                )

            elif method == "DELETE":
                response = session.delete(
                    url, headers=headers, params=query, timeout=timeout
                )

            else:
                raise ValueError(f"Not allowed method: {method}")

            if response.status_code in [429]:
                retry_counter += 1
                time.sleep(retry_interval * retry_counter)
            else:
                break
        except Timeout:
            retry_counter += 1
            if retry_counter > max_retry:
                raise ValueError(f"Timeout\nURL: {url}")

    if response.status_code >= error_threshold:
        raise ValueError(
            f"{response.status_code}\nURL: {response.request.url}\nRequest Header: {response.request.headers}\nRequest Body: {response.request.body}\nResponse Header: {response.headers}\nResponse Body: {response.text}"
        )

    return response


@try_except
def get_switch_keys(account_name):

    path = f"/identity-management/v3/api-clients/self/account-switch-keys"
    query = {
        "search": account_name,
    }

    response = access_api("GET", path, query)
    result = pd.DataFrame.from_dict(response.json())

    return result


@try_except
def list_groups(switch_key, parent_only=False):

    path = "/papi/v1/groups"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }

    response = access_api("GET", path, query, header=header, payload=None)
    result = response.json()
    result = pd.DataFrame.from_dict(result["groups"]["items"])

    if "parentGroupId" in result.columns:
        result = result.sort_values(
            by="parentGroupId", na_position="first", ascending=False
        ).reset_index(drop=True)

    if parent_only:
        try:
            result = result[result["parentGroupId"].isna()]
        except:
            pass

    return result


@try_except
def list_properties(switch_key, contract_id, group_id, network=None):

    path = "/papi/v1/properties"
    query = {
        "accountSwitchKey": switch_key,
        "contractId": contract_id,
        "groupId": group_id,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }

    response = access_api("GET", path, query, header=header, payload=None)
    result = response.json()
    result = pd.DataFrame.from_dict(result["properties"]["items"])

    if network:
        try:
            if network.lower() == "production":
                result = result[result["productionVersion"].notna()].reset_index(
                    drop=True
                )
            elif network.lower() == "staging":
                result = result[result["stagingVersion"].notna()].reset_index(drop=True)
        except:
            pass

    return result


@try_except
def list_hostnames(switch_key, network=None):

    path = "/papi/v1/hostnames"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }

    response = access_api("GET", path, query, header=header, payload=None)
    result = response.json()
    result = pd.DataFrame.from_dict(result["hostnames"]["items"])

    if network:
        if network.lower() == "production":
            result = result.dropna(subset=["productionCnameTo"]).reset_index(drop=True)
        elif network.lower() == "staging":
            result = result.dropna(subset=["stagingCnameTo"]).reset_index(drop=True)

    return result


@try_except
def list_edgehostnames(switch_key, target_edgehostnames=[]):

    path = "/hapi/v1/edge-hostnames"
    query = {
        "accountSwitchKey": switch_key,
    }

    response = access_api("GET", path, query)
    df = pd.DataFrame.from_dict(response.json()["edgeHostnames"])

    for col in df.columns:
        if df[col].dtype == "float64":
            df[col] = df[col].apply(lambda x: int(x) if x.is_integer() else "-")

    for col in ["dnsZone", "recordName"]:
        if col in df.columns:
            df = df[[col] + [c for c in df.columns if c != col]]

    if target_edgehostnames:
        result = []

        for edge_hostname in set(target_edgehostnames):
            try:
                hostname, zone = re.split(r"\.(?=[^.]+\.[^.]+$)", edge_hostname)
            except:
                continue

            df_flt = df[
                (df["recordName"] == hostname.lower()) & (df["dnsZone"] == zone)
            ]
            result.append(df_flt)

        if result:
            result = pd.concat(result)
            result.reset_index(drop=True, inplace=True)
        else:
            return

    else:
        result = df

    return result


@try_except
def list_cpcodes(switch_key, contract_id, group_id):

    path = "/papi/v1/cpcodes"
    query = {
        "accountSwitchKey": switch_key,
        "contractId": contract_id,
        "groupId": group_id,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }

    response = access_api("GET", path, query, header=header)
    result = response.json()
    result = pd.DataFrame.from_dict(result["cpcodes"]["items"])

    return result


@try_except
def list_available_behaviors(switch_key, property_id, version):

    path = f"/papi/v1/properties/{property_id}/versions/{version}/available-behaviors"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }

    response = access_api("GET", path, query, header=header)
    result = response.json()
    result = {
        "product": result["productId"],
        "behaviors": [item["name"] for item in result["behaviors"]["items"]],
    }

    return result


@try_except
def create_property(switch_key, contract_id, group_id, product_id, property_name):

    path = f"/papi/v1/properties"
    query = {
        "accountSwitchKey": switch_key,
        "contractId": contract_id,
        "groupId": group_id,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }
    payload = {
        "productId": product_id,
        "propertyName": property_name,
    }

    response = access_api("POST", path, query, header=header, payload=payload)
    result = response.json()
    result = response.json()["propertyLink"]
    result = int(re.findall(r"\d+$", re.sub(r"[?].*", "", result))[0])
    return result


@try_except
def get_property_info(
    switch_key,
    hostname=None,
    edgehostname=None,
    propertyname=None,
    target_ver=None,
):

    path = "/papi/v1/search/find-by-value"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }
    payload = dict()

    if hostname:
        payload["hostname"] = hostname
    elif edgehostname:
        payload["edgeHostname"] = edgehostname
    elif propertyname:
        payload["propertyName"] = propertyname
    else:
        raise ValueError("One of hostname/edgehostname/propertyname is required")

    response = access_api("POST", path, query, header=header, payload=payload)
    result = response.json()
    result = pd.DataFrame.from_dict(result["versions"]["items"])

    if target_ver:
        if target_ver.lower() == "staging":
            result = result[result["stagingStatus"] == "ACTIVE"]
            result.reset_index(drop=True, inplace=True)
        elif target_ver.lower() == "production":
            result = result[result["productionStatus"] == "ACTIVE"]
            result.reset_index(drop=True, inplace=True)
        elif target_ver.lower() == "latest":
            result = result.sort_values(by="propertyVersion", ascending=False)
            result.reset_index(drop=True, inplace=True)
            result = result.iloc[[0]]

    return result


@try_except
def get_property_activations(switch_key, property_id, network=None, version=None):

    path = f"/papi/v1/properties/{property_id}/activations"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }

    response = access_api("GET", path, query, header=header, payload=None)
    result = response.json()
    result = pd.DataFrame.from_dict(result["activations"]["items"])

    if network:
        if network.upper() in ["PRODUCTION", "STAGING"]:
            result = result[result["network"] == network.upper()]
            result.reset_index(drop=True, inplace=True)

    if version:
        result = result[result["propertyVersion"] == version]
        result.reset_index(drop=True, inplace=True)

    return result


@try_except
def get_property_hostnames(switch_key, property_id, version):

    path = f"/papi/v1/properties/{property_id}/versions/{version}/hostnames"
    query = {
        "accountSwitchKey": switch_key,
        "validateRules": "true",
        "validateMode": "full",
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }

    response = access_api("GET", path, query, header=header)
    result = response.json()
    result = pd.DataFrame.from_dict(result["hostnames"]["items"])

    return result


@try_except
def get_property_json(switch_key, property_id, version):

    path = f"/papi/v1/properties/{property_id}/versions/{version}/rules"
    query = {
        "accountSwitchKey": switch_key,
        "validateRules": "true",
        "validateMode": "full",
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }

    response = access_api("GET", path, query, header=header)
    result = response.json()
    return result


@try_except
def create_edgehostname(
    switch_key,
    contract_id,
    group_id,
    product_id,
    ip_v6,
    hostname,
    edge_hostname_domain,
    enrollment_id=None,
):
    domain_dict = {
        "edgekey.net": "ENHANCED_TLS",
        "edgesuite.net": "STANDARD_TLS",
        "akamaized.net": "SHARED_CERT",
    }

    path = f"/papi/v1/edgehostnames"
    query = {
        "accountSwitchKey": switch_key,
        "contractId": contract_id,
        "groupId": group_id,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }
    payload = {
        "productId": product_id,
        "ipVersionBehavior": "IPV6_COMPLIANCE" if ip_v6 else "IPV4",
        "domainPrefix": re.sub(r"[.]akamaized[.]net$", "", hostname),
        "domainSuffix": edge_hostname_domain,
        "secureNetwork": domain_dict.get(edge_hostname_domain),
    }
    if enrollment_id:
        payload["certEnrollmentId"] = enrollment_id

    response = access_api("POST", path, query, header=header, payload=payload)
    result = response.json()
    result = result.get("edgeHostnameLink")
    if type(result) is str:
        result = int(re.sub(r"\?.+$", "", result).split("/")[-1])

    return result


@try_except
def add_hostname_to_property(
    switch_key,
    contract_id,
    group_id,
    property_id,
    version,
    etag,
    hostname,
    edge_hostname_id,
):

    path = f"/papi/v1/properties/{property_id}/versions/{version}/hostnames"
    query = {
        "accountSwitchKey": switch_key,
        "contractId": contract_id,
        "groupId": group_id,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
        "If-Match": f'"{etag}"',
    }
    payload = {
        "add": [
            {
                "cnameFrom": hostname,
                "edgeHostnameId": str(edge_hostname_id),
            }
        ]
    }

    response = access_api("PATCH", path, query, header=header, payload=payload)
    result = response.json()
    return result


@try_except
def get_edgehostname_info(switch_key, edgehostname):

    edgehostname = re.split(r"\.(?=[^.]+\.[^.]+$)", edgehostname)

    path = f"/hapi/v1/dns-zones/{edgehostname[1]}/edge-hostnames/{edgehostname[0]}"
    query = {
        "accountSwitchKey": switch_key,
    }

    response = access_api("GET", path, query)
    result = pd.DataFrame.from_dict(response.json())

    return result


@try_except
def update_edgehostname_ttl(switch_key, edgehostname, ttl):

    edgehostname = re.split(r"\.(?=[^.]+\.[^.]+$)", edgehostname)

    path = f"/hapi/v1/dns-zones/{edgehostname[1]}/edge-hostnames/{edgehostname[0]}"
    query = {"accountSwitchKey": switch_key, "comments": f"Update TTL to {ttl} sec"}
    header = {
        "accept": "application/json-patch+json",
        "content-type": "application/json-patch+json",
    }
    payload = [{"op": "replace", "path": "/ttl", "value": ttl}]

    response = access_api("PATCH", path, query, header=header, payload=payload)
    result = pd.DataFrame.from_dict(response.json())

    return result


@try_except
def list_bot_categories():

    path = "/appsec/v1/akamai-bot-categories"
    query = None

    response = access_api("GET", path, query, payload=None)
    result = response.json()
    result = pd.DataFrame.from_dict(result["categories"])

    return result


@try_except
def list_bot_detections():

    path = "/appsec/v1/bot-detections"
    query = None

    response = access_api("GET", path, query, payload=None)
    result = response.json()
    result = pd.DataFrame.from_dict(result["detections"])

    return result


@try_except
def list_waf_config(switch_key):

    path = f"/appsec/v1/configs"
    query = {
        "accountSwitchKey": switch_key,
    }

    response = access_api("GET", path, query)
    result = response.json()
    result = pd.DataFrame.from_dict(result["configurations"])

    return result


@try_except
def get_waf_json(switch_key, config_id, version):

    path = f"/appsec/v1/export/configs/{config_id}/versions/{version}"
    query = {
        "accountSwitchKey": switch_key,
    }

    response = access_api("GET", path, query, payload=None)
    result = response.json()

    return result


@try_except
def list_cert_deployments(switch_key, enrollment_id):

    path = f"/cps/v2/enrollments/{enrollment_id}/deployments"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {"accept": "application/vnd.akamai.cps.deployments.v7+json"}

    response = access_api("GET", path, query, header=header)
    result = response.json()
    return result


@try_except
def get_cps_enrollment_id(switch_key, slot=None):

    path = "/cps/v2/enrollments"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {"accept": "application/vnd.akamai.cps.enrollments.v11+json"}

    response = access_api("GET", path, query, header=header)
    result = response.json()
    df = pd.DataFrame.from_dict(result["enrollments"])
    if slot:
        df = df[df["assignedSlots"].apply(lambda x: int(slot) in x)]

    result = df
    return result


@try_except
def get_enrollment(switch_key, enrollment_id):

    path = f"/cps/v2/enrollments/{enrollment_id}"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {
        "accept": "application/vnd.akamai.cps.enrollment.v12+json",
    }

    response = access_api("GET", path, query, header=header)
    result = response.json()
    return result


@try_except
def create_csr(switch_key, contract_id, settings):

    path = f"/cps/v2/enrollments"
    query = {
        "accountSwitchKey": switch_key,
        "contractId": contract_id,
    }
    header = {
        "accept": "application/vnd.akamai.cps.enrollment-status.v1+json",
        "content-type": "application/vnd.akamai.cps.enrollment.v12+json",
    }
    payload = settings

    response = access_api("POST", path, query, header=header, payload=payload)
    result = response.json()
    result = result.get("enrollment")
    if type(result) is str:
        result = int(re.sub(r"\?.+$", "", result).split("/")[-1])

        return result


@try_except
def update_cps_deployment(switch_key, enrollment_id, cps_config, cancel_pending=False):

    path = f"/cps/v2/enrollments/{enrollment_id}"
    query = {
        "accountSwitchKey": switch_key,
        "allow-cancel-pending-changes": cancel_pending,
    }
    header = {
        "accept": "application/vnd.akamai.cps.enrollment-status.v1+json",
        "content-type": "application/vnd.akamai.cps.enrollment.v12+json",
    }
    payload = cps_config

    response = access_api("PUT", path, query, header=header, payload=payload)
    result = response.json()
    result = result.get("changes")
    return result


@try_except
def get_cps_change(switch_key, enrollment_id, change_id):

    path = f"/cps/v2/enrollments/{enrollment_id}/changes/{change_id}/input/info/change-management-info"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {
        "accept": "application/vnd.akamai.cps.deployment.v8+json",
    }

    response = access_api("GET", path, query, header=header)
    result = response.json()
    return result


@try_except
def create_property_version(switch_key, property_id, base_version):

    path = f"/papi/v1/properties/{property_id}/versions"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {"PAPI-Use-Prefixes": "false"}
    payload = {
        "createFromVersion": str(base_version),
    }

    response = access_api("POST", path, query, header=header, payload=payload)
    result = response.json()["versionLink"]
    result = int(re.findall(r"\d+$", re.sub(r"[?].*", "", result))[0])
    return result


@try_except
def update_property(switch_key, property_id, version, rules, etag, note, dryrun=False):

    path = f"/papi/v1/properties/{property_id}/versions/{version}/rules"
    query = {
        "accountSwitchKey": switch_key,
    }
    if dryrun:
        query["dryRun"] = "true"
    header = {
        "PAPI-Use-Prefixes": "false",
        "If-Match": f'"{etag}"',
    }
    payload = {"rules": rules, "comments": str(note)}

    response = access_api("PUT", path, query, header=header, payload=payload)
    result = response.json()
    return result


@try_except
def update_behaviors(rules, behavior_name, option_name, option_value=None):

    # Update option
    if option_name is None:
        rules["behaviors"] = [
            behavior
            for behavior in rules["behaviors"]
            if behavior["name"] != behavior_name
        ]
    else:
        for behavior in rules["behaviors"]:
            if behavior["name"] == behavior_name:
                behavior["options"][option_name] = option_value

    # Check Child rules
    for child_rule in rules["children"]:
        update_behaviors(child_rule, behavior_name, option_name, option_value)

    return rules


@try_except
def activate_staging(
    switch_key, property_id, version, note, emails=[], fast_act=False, validation=True
):

    path = f"/papi/v1/properties/{property_id}/activations"
    query = {
        "accountSwitchKey": switch_key,
        "validateRules": validation,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }
    payload = {
        "propertyVersion": str(version),
        "network": "STAGING",
        "note": str(note),
        "notifyEmails": emails,
        "fastPush": "true" if fast_act else "false",
        "acknowledgeAllWarnings": "true",
    }

    response = access_api("POST", path, query, header=header, payload=payload)
    result = response.json()
    return result


@try_except
def create_cpcode(switch_key, contract_id, group_id, product_id, cpcode_name):

    path = f"/papi/v1/cpcodes"
    query = {
        "accountSwitchKey": switch_key,
        "contractId": contract_id,
        "groupId": group_id,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }
    payload = {
        "productId": product_id,
        "cpcodeName": cpcode_name,
    }

    response = access_api("POST", path, query, header=header, payload=payload)
    result = response.json()
    result = result.get("cpcodeLink")
    if type(result) is str:
        result = int(re.sub(r"\?.+$", "", result).split("/")[-1])

    return result


@try_except
def get_cpcode(switch_key, contract_id, group_id, cpcode):

    path = f"/papi/v1/cpcodes/cpc_{cpcode}"
    query = {
        "accountSwitchKey": switch_key,
        "contractId": contract_id,
        "groupId": group_id,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }

    response = access_api("GET", path, query, header=header)
    result = response.json()
    result = pd.DataFrame.from_dict(result["cpcodes"]["items"])

    return result


@try_except
def list_custom_behaviors(switch_key):

    path = f"/papi/v1/custom-behaviors"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }

    response = access_api("GET", path, query, header=header)
    result = response.json()
    result = pd.DataFrame.from_dict(result["customBehaviors"]["items"])
    return result


@try_except
def list_custom_overrides(switch_key):

    path = f"/papi/v1/custom-overrides"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {
        "PAPI-Use-Prefixes": "false",
    }

    response = access_api("GET", path, query, header=header)
    result = response.json()
    result = pd.DataFrame.from_dict(result["customOverrides"]["items"])
    return result


@try_except
def get_dns_records(switch_key, zone):

    path = f"/config-dns/v2/zones/{zone}/zone-file"
    query = {
        "accountSwitchKey": switch_key,
    }
    header = {"accept": "text/dns"}

    response = access_api("GET", path, query, header=header)
    result = response.text
    return result


@try_except
def get_dns_zone_info(switch_key, zone_name):

    path = f"/config-dns/v2/zones/{zone_name}"
    query = {
        "accountSwitchKey": switch_key,
    }

    response = access_api("GET", path, query)
    result = response.json()
    return result


@try_except
def get_dns_record_sets(
    switch_key,
    zone_name,
    record_name,
    record_type,
):

    path = f"/config-dns/v2/zones/{zone_name}/recordsets"
    query = {
        "accountSwitchKey": switch_key,
        "types": record_type,
        "search": record_name,
    }

    response = access_api("GET", path, query)
    result = response.json()
    result = pd.DataFrame.from_dict(result["recordsets"])
    return result


@try_except
def create_dns_zone(
    switch_key,
    contract_id,
    group_id,
    zone_name,
    zone_type="PRIMARY",
    master_servers=[],
    dns_sec=False,
    algorithm=None,
    tsig=None,
):

    path = f"/config-dns/v2/zones"
    query = {
        "accountSwitchKey": switch_key,
        "contractId": contract_id,
        "groupId": group_id,
    }
    payload = {
        "type": zone_type,
        "zone": zone_name,
    }

    if zone_type == "SECONDARY":
        payload["masters"] = master_servers

    if dns_sec:
        payload["signAndServe"] = dns_sec
        payload["signAndServeAlgorithm"] = algorithm
        payload["tsigKey"] = tsig

    response = access_api("POST", path, query, payload=payload)
    result = response.json()
    return result


@try_except
def create_dns_record_set(
    switch_key,
    zone_name,
    record_name,
    record_type,
    record_data,
    ttl=300,
    replace=False,
):

    path = f"/config-dns/v2/zones/{zone_name}/names/{record_name}/types/{record_type}"
    query = {
        "accountSwitchKey": switch_key,
    }
    payload = {
        "name": record_name,
        "type": record_type,
        "rdata": record_data,
        "ttl": int(ttl),
    }

    response = access_api("PUT" if replace else "POST", path, query, payload=payload)
    result = response.json()
    return result


@try_except
def create_test_suite(switch_key, suite_name):

    path = "/test-management/v3/functional/test-suites"
    query = {"accountSwitchKey": switch_key}
    payload = {
        "testSuiteName": f"{suite_name} - {datetime.now().strftime('%Y/%m/%d %H:%M:%S')}",
        "testSuiteDescription": "Created by API_v3",
        "isLocked": False,
        "isStateful": False,
    }

    response = access_api("POST", path, query, payload=payload)
    result = response.json()
    return result["testSuiteId"]


@try_except
def include_test_cases(switch_key, suite_id, test_cases):

    path = f"/test-management/v3/functional/test-suites/{suite_id}/test-cases"
    query = {"accountSwitchKey": switch_key}
    payload = test_cases

    response = access_api("POST", path, query, payload=payload)
    result = response.json()
    return result


@try_except
def include_test_variable(switch_key, suite_id, test_variables):

    path = f"/test-management/v3/functional/test-suites/{suite_id}/variables"
    query = {"accountSwitchKey": switch_key}
    payload = [
        {"variableName": k, "variableValue": v} for k, v in test_variables.items()
    ]

    response = access_api("POST", path, query, payload=payload)
    result = response.json()
    return result


@try_except
def list_datastreams(switch_key):

    path = "/datastream-config-api/v2/log/streams"
    query = {"accountSwitchKey": switch_key}

    response = access_api("GET", path, query)
    result = pd.DataFrame.from_dict(response.json())
    return result


@try_except
def get_datastream(switch_key, stream_id):

    path = f"/datastream-config-api/v2/log/streams/{stream_id}"
    query = {"accountSwitchKey": switch_key}

    response = access_api("GET", path, query)
    result = response.json()
    return result


@try_except
def create_datastream(switch_key, config):

    path = "/datastream-config-api/v2/log/streams"
    query = {"accountSwitchKey": switch_key}
    payload = config

    response = access_api("POST", path, query, payload=payload)
    result = response.json()
    return result


@try_except
def list_available_dataset_fields(switch_key):

    path = "/datastream-config-api/v2/log/datasets-fields"
    query = {"accountSwitchKey": switch_key}

    response = access_api("GET", path, query)
    result = response.json()
    result = {
        d["datasetFieldName"]: d["datasetFieldId"] for d in result["datasetFields"]
    }
    return result


@try_except
def get_traffic_stats(
    switch_key, property_id, start, end, dimensions=None, format_value=True
):

    if dimensions:
        if type(dimensions) is str:
            dimensions = [dimensions]
        else:
            try:
                dimensions = list(dimensions)
            except:
                dimensions = [dimensions]
    else:
        dimensions = []

    path = "/reporting-api/v2/reports/delivery/traffic/current/data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
    }
    payload = {
        "dimensions": list(dimensions) if dimensions else [],
        "metrics": [
            "edgeHitsSum",
            "originHitsSum",
            "offloadedHitsPercentage",
            "edgeBytesSum",
            "originBytesSum",
            "offloadedBytesPercentage",
        ],
        "filters": [
            {
                "dimensionName": "arlId",
                "operator": "IN_LIST",
                "expressions": [property_id],
            },
        ],
        "sortBys": (
            [{"name": dimension, "sortOrder": "ASCENDING"} for dimension in dimensions]
            if dimensions
            else []
        ),
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])

    if format_value:
        for col in result.columns:
            try:
                if "offload" in col.lower():
                    result[col] = result[col].apply(lambda x: f"{x:.2f}%")
                elif "hits" in col.lower():
                    result[col] = result[col].apply(lambda x: "{:,}".format(int(x)))
                elif "bytes" in col.lower():
                    result[col] = result[col].apply(convert_bytes)
                elif col in ["time5minutes", "time1hour", "time1day"]:
                    result[col] = pd.to_datetime(result[col], unit="s", utc=True)
                    result[col] = result[col].dt.strftime("%Y-%m-%dT%H:%M:%SZ")
            except:
                pass

    return result


@try_except
def report_on_browser(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/hits-by-browser/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "DAY",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["successfulHits", "successfulHitsPercent"],
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])
    return result


@try_except
def report_on_os(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/hits-by-os/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "DAY",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["successfulHits", "successfulHitsPercent"],
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])
    return result


@try_except
def report_on_geo(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/enhancedtraffic-by-country/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "DAY",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["country", "edgeHits"],
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])
    return result


@try_except
def report_on_url_hits(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/urlhits-by-url/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "DAY",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["allEdgeHits", "allOriginHits", "allHitsOffload"],
        "limit": 500,
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])
    return result


@try_except
def report_on_url_bytes(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/urlbytes-by-url/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "DAY",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["allEdgeBytes", "allOriginBytes", "allBytesOffload"],
        "limit": 500,
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])
    return result


@try_except
def report_on_url_4xx(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/url4XXresponses-by-url/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "DAY",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["404EdgeHits", "4XXOtherEdgeHits"],
        "limit": 500,
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])
    return result


@try_except
def report_on_url_3xx(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/url3XXresponses-by-url/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "DAY",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["302EdgeHits", "304EdgeHits", "3XXOtherEdgeHits"],
        "limit": 500,
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])
    return result


@try_except
def report_on_url_2xx(switch_key, cpcode, start, end, hostnames=None):

    path = "/reporting-api/v1/reports/url2XXresponses-by-url/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "DAY",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["200EdgeHits", "206EdgeHits", "2XXOtherEdgeHits"],
        "limit": 500,
    }

    if hostnames:
        payload["filters"] = {
            "url_start_with": hostnames if type(hostnames) is list else list(hostnames)
        }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])
    return result


@try_except
def report_on_response_code(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/traffic-by-response/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "DAY",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "filters": {"response_status": ["success", "error"]},
        "metrics": ["edgeHits", "edgeHitsPercent", "originHits", "originHitsPercent"],
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])
    return result


@try_except
def report_on_hits_5xx(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/traffic-by-timeandresponseclass/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "HOUR",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "filters": {"response_status": ["error"], "response_class": ["5xx"]},
        "metrics": ["edgeHitsPerSecond", "originHitsPerSecond"],
    }

    response = access_api("POST", path, query, payload=payload)
    df = pd.DataFrame.from_dict(response.json()["data"])
    result = pd.concat(
        [df["startdatetime"], pd.DataFrame([d[0] for d in df["data"]])], axis=1
    )
    return result


@try_except
def report_on_hits_0(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/traffic-by-timeandresponseclass/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "HOUR",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "filters": {"response_status": ["error"], "response_class": ["0xx"]},
        "metrics": ["edgeHitsPerSecond", "originHitsPerSecond"],
    }

    response = access_api("POST", path, query, payload=payload)
    df = pd.DataFrame.from_dict(response.json()["data"])
    result = pd.concat(
        [df["startdatetime"], pd.DataFrame([d[0] for d in df["data"]])], axis=1
    )
    return result


@try_except
def report_on_hits(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/hits-by-time/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "HOUR",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["edgeHitsPerSecond", "originHitsPerSecond", "hitsOffload"],
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])
    return result


@try_except
def report_on_bytes(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/bytes-by-time/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "HOUR",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["edgeBitsPerSecond", "originBitsPerSecond", "bytesOffload"],
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["data"])
    return result


@try_except
def report_on_hits_total(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/hits-by-time/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "HOUR",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["edgeHitsTotal", "originHitsTotal"],
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["summaryStatistics"])

    edge_total = float(result.at["value", "edgeHitsTotal"])
    origin_total = float(result.at["value", "originHitsTotal"])
    result.at["value", "offloadHitsTotal"] = (
        f"{max(0, round(100 * (edge_total - origin_total) / edge_total, 2)):.2f}%"
        if edge_total > 0
        else "-"
    )

    for col in result.columns.values:
        if ("hits" in col.lower()) and ("offload" not in col.lower()):
            try:
                result.at["value", col] = "{:,}".format(int(result.at["value", col]))
            except:
                pass

    return result


@try_except
def report_on_bytes_total(switch_key, cpcode, start, end):

    path = "/reporting-api/v1/reports/bytes-by-time/versions/1/report-data"
    query = {
        "accountSwitchKey": switch_key,
        "start": start.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "end": end.strftime("%Y-%m-%dT%H:%M:%SZ"),
        "interval": "HOUR",
    }
    payload = {
        "objectType": "cpcode",
        "objectIds": list(set(cpcode)) if type(cpcode) is list else list(cpcode),
        "metrics": ["edgeBytesTotal", "originBytesTotal"],
    }

    response = access_api("POST", path, query, payload=payload)
    result = pd.DataFrame.from_dict(response.json()["summaryStatistics"])

    edge_total = float(result.at["value", "edgeBytesTotal"])
    origin_total = float(result.at["value", "originBytesTotal"])
    result.at["value", "offloadBytesTotal"] = (
        f"{max(0, round(100 * (edge_total - origin_total) / edge_total, 2)):.2f}%"
        if edge_total > 0
        else "-"
    )

    for col in result.columns.values:
        if ("bytes" in col.lower()) and ("offload" not in col.lower()):
            result.at["value", col] = convert_bytes(result.at["value", col])

    return result


@try_except
def create_ppt_slide(presentation, layout, slide_title="Title"):

    slide_layout = presentation.slide_layouts[layout]
    slide = presentation.slides.add_slide(slide_layout)

    if slide.shapes.title:
        slide.shapes.title.text = slide_title

    else:
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.text = slide_title
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
                break

    return slide


@try_except
def add_table_to_slide(
    slide,
    df,
    left=0.5,
    top=1,
    width=5,
    height=0.5,
    font_size=6,
    include_index=False,
    auto_adjust=False,
):

    rows, cols = df.shape
    col_offset = 1 if include_index else 0
    top = Inches(top)
    left = Inches(left)
    width = Inches(width)
    height = Inches(height)
    table = slide.shapes.add_table(
        rows + 1, (cols + 1) if include_index else cols, left, top, width, height
    ).table

    fill_color = RGBColor(13, 154, 220)
    font_name = "Arial"
    font_size = Pt(font_size)
    alignment = PP_ALIGN.RIGHT

    # Set index
    if include_index:
        cell = table.cell(0, 0)
        table.columns[0].width = Inches(0.3)
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color

        text_frame = cell.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = font_name
            paragraph.font.size = font_size

        for i, index in enumerate(df.index.values):
            cell = table.cell(i + 1, 0)
            cell.fill.solid()

            if i % 2:
                cell.fill.fore_color.rgb = RGBColor(231, 231, 231)
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            text_frame = cell.text_frame
            text_frame.text = str(index)
            text_frame.word_wrap = False

            for paragraph in text_frame.paragraphs:
                paragraph.font.name = font_name
                paragraph.font.size = font_size
                paragraph.alignment = PP_ALIGN.LEFT

    # Set column names
    for j, column in enumerate(df.columns):
        cell = table.cell(0, j + col_offset)
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color

        text_frame = cell.text_frame
        text_frame.text = str(column) if "percent" not in str(column).lower() else "%"

        for paragraph in text_frame.paragraphs:
            paragraph.font.name = font_name
            paragraph.font.size = font_size
            paragraph.alignment = PP_ALIGN.LEFT if j == 0 else alignment

    # Set values
    for i, _ in enumerate(df.index.values):
        for j, value in enumerate(df.iloc[i]):
            cell = table.cell(i + 1, j + col_offset)
            cell.fill.solid()

            if i % 2:
                cell.fill.fore_color.rgb = RGBColor(231, 231, 231)
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            if type(value) is float:
                value = int(value) if value.is_integer() else value

            text_frame = cell.text_frame
            text_frame.text = str(value)
            text_frame.word_wrap = False

            for paragraph in text_frame.paragraphs:
                paragraph.font.name = font_name
                paragraph.font.size = font_size

            if j == 0:
                text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            else:
                text_frame.paragraphs[0].alignment = alignment

    # Adjust column width
    if auto_adjust:
        for j, column in enumerate(df.columns):
            df[column] = df[column].astype(str)
            len_max = max(
                df[column].apply(len).max(),
                len(table.cell(0, j + col_offset).text_frame.text),
            )
            table.columns[j + col_offset].width = min(
                max(Inches(0.1) * len_max, Inches(0.4)), Inches(5)
            )


@try_except
def add_text_to_slide(slide, text, font_size=16):

    font_name = "Arial"
    font_size = Pt(font_size)

    for shape in slide.shapes:

        if shape.has_text_frame and not shape.text_frame.text:
            shape.text_frame.text = text

            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.name = font_name
                paragraph.font.size = font_size

            break


@try_except
def add_image_to_slide(slide, file_name, left, top, width=None, height=None):

    left = Inches(left)
    top = Inches(top)
    width = Inches(width) if width is not None else width
    height = Inches(height) if height is not None else height

    slide.shapes.add_picture(file_name, left, top, width, height)